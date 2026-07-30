import fitz
from pptx import Presentation
from pptx.util import Pt
from pathlib import Path


#name = "大三島工場製造資料20260713追加_包装2班.pdf"
#pdf_path = rf"C:\Users\skacyba\Downloads\{name}"
#pptx_path = rf"{name}.pptx"

template = r"template.pptx"

folder = Path(r".")
prs = Presentation(template)
        
for pdf_path in folder.glob("*.pdf"):
    pptx_path = rf"{pdf_path}.pptx"
    doc = fitz.open(pdf_path)
    blank = prs.slide_layouts[6]

    for i, page in enumerate(doc):
        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
        img_path = f"page_{i+1}.png"
        pix.save(img_path)

        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(img_path, 0, 0, width=prs.slide_width, height=prs.slide_height)

        filename_box = slide.shapes.add_textbox(0, 0, prs.slide_width, Pt(12))
        filename_frame = filename_box.text_frame
        filename_frame.margin_left = 0
        filename_frame.margin_right = 0
        filename_frame.margin_top = 0
        filename_frame.margin_bottom = 0

        filename_paragraph = filename_frame.paragraphs[0]
        filename_run = filename_paragraph.add_run()
        filename_run.text = pdf_path.name
        filename_run.font.name = "メイリオ"
        filename_run.font.size = Pt(9)

prs.save(pptx_path)

